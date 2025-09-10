import os
import tempfile
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

def parse_docx(file_content):
    """解析Word文档"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
        temp_file.write(file_content)
        temp_file_path = temp_file.name
        
    doc = Document(temp_file_path)
    content = []
    for para in doc.paragraphs:
        if para.text.strip():
            content.append({
                'type': 'paragraph',
                'content': para.text,
                'position': f'paragraph_{len(content)}'
            })
    
    # 清理临时文件
    os.unlink(temp_file_path)
    return {
        'document_type': 'docx',
        'content': content
    }

def parse_xlsx(file_content):
    """解析Excel文档"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
        temp_file.write(file_content)
        temp_file_path = temp_file.name
        
    workbook = load_workbook(temp_file_path, read_only=True)
    content = []
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
            for col_idx, cell_value in enumerate(row):
                if cell_value is not None and str(cell_value).strip():
                    content.append({
                        'type': 'cell',
                        'content': str(cell_value),
                        'position': f'sheet_{sheet_name}_row_{row_idx}_col_{col_idx}'
                    })
    
    workbook.close()
    os.unlink(temp_file_path)
    return {
        'document_type': 'xlsx',
        'content': content
    }

def parse_pptx(file_content):
    """解析PowerPoint文档"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
        temp_file.write(file_content)
        temp_file_path = temp_file.name
        
    prs = Presentation(temp_file_path)
    content = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                content.append({
                    'type': 'slide_shape',
                    'content': shape.text,
                    'position': f'slide_{slide_idx}_shape_{shape_idx}'
                })
    
    os.unlink(temp_file_path)
    return {
        'document_type': 'pptx',
        'content': content
    }

def main(inputs):
    # 获取上传的文件内容和文件名
    file_content = inputs.get('file_content')
    file_name = inputs.get('file_name', '')
    
    # 根据文件扩展名选择相应的解析方法
    if file_name.endswith('.docx'):
        return parse_docx(file_content)
    elif file_name.endswith('.xlsx'):
        return parse_xlsx(file_content)
    elif file_name.endswith('.pptx'):
        return parse_pptx(file_content)
    else:
        raise ValueError(f"不支持的文件格式: {file_name}")
