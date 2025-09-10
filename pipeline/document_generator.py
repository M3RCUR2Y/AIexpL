import os
import tempfile
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import json

def generate_docx(original_content, modified_content):
    """生成修改后的Word文档"""
    doc = Document()
    modified_data = json.loads(modified_content)
    
    # 简单重建文档（实际应用中可能需要更复杂的格式保留逻辑）
    for item in modified_data:
        if item['type'] == 'paragraph':
            doc.add_paragraph(item['content'])
    
    # 保存到临时文件
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
        doc.save(temp_file)
        temp_file_path = temp_file.name
    
    # 读取文件内容
    with open(temp_file_path, 'rb') as f:
        file_content = f.read()
    
    # 清理临时文件
    os.unlink(temp_file_path)
    return file_content

def generate_xlsx(original_content, modified_content):
    """生成修改后的Excel文档"""
    wb = Workbook()
    ws = wb.active
    modified_data = json.loads(modified_content)
    
    # 构建单元格位置到内容的映射
    cell_map = {}
    for item in modified_data:
        if item['type'] == 'cell':
            parts = item['position'].split('_')
            sheet_name = parts[1]
            row = int(parts[3]) + 1  # Excel行号从1开始
            col = int(parts[5]) + 1  # Excel列号从1开始
            cell_map[(sheet_name, row, col)] = item['content']
    
    # 确保工作表存在
    if 'sheet_1' not in wb.sheetnames:
        wb.create_sheet('1')
    ws = wb['1']
    
    # 填充数据
    for (sheet, row, col), content in cell_map.items():
        if sheet == '1':  # 简化处理，仅处理第一个工作表
            ws.cell(row=row, column=col, value=content)
    
    # 保存到临时文件
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
        wb.save(temp_file)
        temp_file_path = temp_file.name
    
    # 读取文件内容
    with open(temp_file_path, 'rb') as f:
        file_content = f.read()
    
    # 清理临时文件
    os.unlink(temp_file_path)
    return file_content

def generate_pptx(original_content, modified_content):
    """生成修改后的PowerPoint文档"""
    prs = Presentation()
    modified_data = json.loads(modified_content)
    
    # 简单重建文档
    slide_content_map = {}
    for item in modified_data:
        if item['type'] == 'slide_shape':
            parts = item['position'].split('_')
            slide_idx = int(parts[1])
            shape_idx = int(parts[3])
            
            if slide_idx not in slide_content_map:
                slide_content_map[slide_idx] = []
            slide_content_map[slide_idx].append((shape_idx, item['content']))
    
    # 创建幻灯片并添加内容
    for slide_idx in sorted(slide_content_map.keys()):
        if slide_idx == 0:
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片
        else:
            slide_layout = prs.slide_layouts[1]  # 标题和内容
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加内容
        for shape_idx, content in slide_content_map[slide_idx]:
            if shape_idx < len(slide.shapes) and hasattr(slide.shapes[shape_idx], "text"):
                slide.shapes[shape_idx].text = content
    
    # 保存到临时文件
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
        prs.save(temp_file)
        temp_file_path = temp_file.name
    
    # 读取文件内容
    with open(temp_file_path, 'rb') as f:
        file_content = f.read()
    
    # 清理临时文件
    os.unlink(temp_file_path)
    return file_content

def main(inputs):
    original_content = inputs.get('original_content')
    modified_content = inputs.get('modified_content')
    document_type = inputs.get('document_type')
    
    if document_type == 'docx':
        return generate_docx(original_content, modified_content)
    elif document_type == 'xlsx':
        return generate_xlsx(original_content, modified_content)
    elif document_type == 'pptx':
        return generate_pptx(original_content, modified_content)
    else:
        raise ValueError(f"不支持的文档类型: {document_type}")
